import pptx

prs = pptx.Presentation("powerpoints/test.pptx")


def print_all_text(presentation):
    """print all runs in a presentation"""
    # text_runs will be populated with a list of strings,
    # one for each text run in presentation
    text_runs = []

    for slide in presentation.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.text == "<source>":
                    print("p:", paragraph.text)
                for run in paragraph.runs:
                    text_runs.append(run.font.name)

    print(text_runs[-27:-15])

print_all_text(prs)
# prs.save('test.pptx')