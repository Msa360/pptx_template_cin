from utils import add_countries, add_subtitle, add_title, clean_up_shapes
import pptx
from docx_parser import full_tree
import sys

if __name__ == "__main__":
    print(sys.argv[1], sys.argv[2])
    in_file = sys.argv[1]
    out_file = sys.argv[2]
    state_dict = full_tree(in_file)
    prs = pptx.Presentation("powerpoints/gabarit_mod.pptx")
    add_title(prs.slides[0], state_dict["title"])
    add_subtitle(prs.slides[0], state_dict["subtitle"])
    top = 100000
    bottom = add_countries(prs, prs.slides[-1], state_dict["countries"], top)
    clean_up_shapes(prs, "<banner>") # deletes the template shapes
    prs.save(out_file)
