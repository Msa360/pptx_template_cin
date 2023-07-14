from tkinter import Tk, Button, Label, Entry, filedialog
import tkinter as tk
import webbrowser
import os, sys
from ui.loading import StatusScreen

def resource_path(relative_path):
    """ Get absolute path to resource, works for dev and for PyInstaller """
    try:
        # PyInstaller creates a temp folder and stores path in _MEIPASS
        base_path = sys._MEIPASS
    except Exception:
        base_path = os.path.abspath(".")

    return os.path.join(base_path, relative_path)



class MainWindow(Tk):
    def __init__(self):
        """
        This is the main window of the app.
        """
        super().__init__()

        self.title("Word-2-Powerpoint")
        self.label_bienvenue = Label(text="Veuiller chosir un fichier word a transformer")
        self.label_bienvenue.grid(row=0, column=0, padx=10, pady=10)

        self.btn_input_file = Button(text="Choisir un fichier", width=20, command=self.choose_input_file)
        self.btn_input_file.grid(row=1, column=0, padx=10, pady=10)

        self.file_entry_text = tk.StringVar()
        self.file_entry = Entry(self, width=30, textvariable=self.file_entry_text)
        self.file_entry.grid(row=2, column=0, padx=10, pady=0)

        self.btn_output_file = Button(text="Choisir le dossier sortant", width=20, command=self.choose_output_dir)
        self.btn_output_file.grid(row=3, column=0, padx=10, pady=10)

        self.dir_entry_text = tk.StringVar()
        self.dir_entry = Entry(self, width=30, textvariable=self.dir_entry_text)
        self.dir_entry.grid(row=4, column=0, padx=10, pady=0)

        self.btn_make_pptx = Button(text="Go", width=20, command=self.make_pptx)
        self.btn_make_pptx.grid(row=5, column=0, padx=10, pady=10)


        self.help_link = Label(text="Aide", font=('Helveticabold', 15), fg="blue")
        self.help_link.grid(row=6, column=0, padx=10, pady=10)
        self.help_link.bind("<Button-1>", lambda e: webbrowser.open_new_tab("https://github.com/Msa360/pptx_template_cin/blob/master/docs.md"))

    def choose_input_file(self):
        filename = filedialog.askopenfilename(parent=self)
        if filename != "":
            self.file_entry_text.set(filename)

    def choose_output_dir(self):
        dir = filedialog.askdirectory(parent=self, mustexist=True)
        if dir != "":
            self.dir_entry_text.set(dir)

    def make_output_file(self):
        output_file = os.path.join(self.dir_entry_text.get(), os.path.basename(self.file_entry_text.get()).rsplit(".", 1)[0]+".pptx")
        i = 1
        while os.path.exists(output_file):
            if i == 1:
                output_file = output_file[:-5] + str(i) + ".pptx"
            else:
                output_file = output_file[:-6] + str(i) + ".pptx"
            i += 1
        return output_file

    def make_pptx(self):
        # create a loading window
        # loading_screen = LoadingScreen(self)
        # call pptx creator todo: make this code better
        try:
            import pptx
            import word2pptx as wpx

            state_dict = wpx.word_tree(self.file_entry_text.get())
            prs = pptx.Presentation(resource_path(os.path.join("powerpoints", "gabarit_v1.pptx")))
            top = 3_000_000 # decided arbitrarily
            wpx.add_title(prs.slides[0], state_dict["title"])
            wpx.add_subtitle(prs.slides[0], state_dict["subtitle"])
            slide, bottom = wpx.add_intro(prs, prs.slides[0], state_dict["intro"], top)
            bottom = wpx.add_countries(prs, slide, state_dict["countries"], bottom)
            wpx.add_sources(prs, prs.slides[1], state_dict["sources"], 0)
            # deletes the template shapes
            wpx.clean_up_shapes(prs, "<banner>") 
            wpx.clean_up_shapes(prs, "<source_banner>")
            # clean_up_shapes(prs, "<source>") # issues with because it deletes the whole shape
            output_file = self.make_output_file()
            prs.save(output_file)
            # end call pptx creator
        except Exception as e:
            loading_screen = StatusScreen(self, str(e))
        else:
            loading_screen = StatusScreen(self, "Succès!")
        self.wait_window(loading_screen)