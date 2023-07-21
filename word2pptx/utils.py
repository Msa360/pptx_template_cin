import copy
import pptx
from pptx.shapes.autoshape import Shape
from pptx.enum.text import PP_ALIGN
from typing import List, Dict

# GLOBAL VARS
SLIDE_LAYOUT = 4 # style of slide
MAX_CHARS_PER_LINE = 97 # approximate, experimentally determined
BANNER_HEIGHT = 400_000 # defined arbitrarily to ressemble canvas
BANNER_MARKUP = "<banner>"
SOURCE_BANNER_MARKUP = "<source_banner>"
SOURCE_MARKUP = "<source>"
SUBTITLE_LMARGIN = 138545
AUTHOR_MARKUP = "<author>"



def chunck_text(text: str, lenght: int):
    """chuncks the text into lines that are smaller or equal than 'lenght' arg"""
    text = list(text)
    text_list = list(text)
    # for i in range(len(text) // lenght):  # iter over lines
    pos = 0
    while len(text_list) / lenght > 1:
        for j in reversed(range(0, lenght+1)):    # iter over chars
            if text_list[j] == " ":
                pos += j
                text[pos] = "\n"
                pos += 1
                text_list = text_list[j+1:]
                break
            if j == 0:
                raise Exception(f"Cannot chunck {''.join(text_list[:lenght])}|{text_list[lenght]}, it would break the word at '|'")
    return ''.join(text)

def remove_newlines(text: str):
    new_text = list(text)
    for i in range(len(text)):
        if 0 < i < len(text) - 1:
            if text[i] == '\n':
                if text[i-1] == ' ' and text[i+1] == ' ':
                    new_text[i] = ''
                    new_text[i-1] = ''
                elif text[i-1] == ' ' and text[i+1] != ' ':
                    new_text[i] = ''
                elif text[i-1] != ' ' and text[i+1] == ' ':
                    new_text[i] = ''
                elif text[i-1] != ' ' and text[i+1] != ' ':
                    new_text[i] = ' '
        else:
            if text[i] == '\n':
                new_text[i] = ''
    return new_text

# a very important algo
def textbox_height_estimate(text: str):
    """
    The goal is to estimate the height of the textbox to fit all the characters
    this is needed since pptx api can't do it because it would need a rendering engine
    """
    n_chars = len(text)
    return (max(n_chars - 100, 0)) * 0.000183 + 0.03


def is_slide_full(presentation, shape_top, shape_height):
    """check if the content will fit in the slide"""
    # Get the last slide in the presentation
    slide = presentation.slides[-1]
    shapes = list(slide.shapes)
    if len(shapes)==0 or presentation.slide_height > shape_top + shape_height:
        return False
    else:
        return True
    

def clone_shape(shape, left, top, width, height):
        # gotta be careful since the shape_id is read-only, it keeps the same id so the id isn't unique anymore
        # gotta be careful, the style can change because of hierarchies
        """
        Add a duplicate of `shape` to the same slide.
        Once duplicated, id is obselete, use `shape.element` to have uniqueness.
        """
        shape_obj = shape.element
        sp_tree = shape_obj.getparent()
        new_sp = copy.deepcopy(shape_obj)
        sp_tree.append(new_sp)
        new_shape = Shape(new_sp, None) #(GraphicFrame(new_sp, None) if isinstance(shape, GraphicFrame) else Shape(new_sp, None))
        new_shape.left = int(left)
        new_shape.top = int(top)
        new_shape.width = int(width)
        new_shape.height = int(height)
        return new_shape

def delete_shape(shape, slide):
    slide.shapes._spTree.remove(shape._element)

def move_slide(presentation, old_index, new_index):
    slides = list(presentation.slides._sldIdLst)
    presentation.slides._sldIdLst.remove(slides[old_index])
    presentation.slides._sldIdLst.insert(new_index, slides[old_index])

def find_template_shape(slide, markup: str):
    """
    markup is the template placeholder

    raise error in markup isn't found

    returns a generator of shapes object
    """
    used = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text == markup:
                yield shape
                used += 1
    if used == 0:
        raise Exception(f"template shape wasn't found in {slide} with {markup}")

def clean_up_shapes(presentation, markup: str):
    """deletes template shapes"""
    for slide in presentation.slides:
        try:
            for shape in find_template_shape(slide, markup):
                delete_shape(shape, slide)
        except:
            pass


def add_title(slide, title: str):
    TITLE_MARKUP = "<title>"
    title_shape = next(find_template_shape(slide, TITLE_MARKUP))
    title_shape.text_frame.text = title


def add_subtitle(slide, subtitle: str):
    SUBTITLE_MARKUP = "<subtitle>"
    subtitle_shape = next(find_template_shape(slide, SUBTITLE_MARKUP))
    subtitle_shape.text_frame.text = subtitle

def add_intro(presentation, slide, intro: List[Dict], top: int):

    bottom = top # at the beginning
    for bodypart in intro:
        kv = tuple(bodypart.items())
        if len(kv) == 1:
            k, v = kv[0]
        else:
            raise Exception(f"Error while reading {bodypart} dict from article tree dict")
        
        if k == 'subtitle':
            # add subtitle
            subtitle_height = presentation.slide_height * textbox_height_estimate(v)
            subtitle = slide.shapes.add_textbox(SUBTITLE_LMARGIN, bottom, presentation.slide_width, subtitle_height)
            subtitle_text = subtitle.text_frame.paragraphs[0]
            subtitle_text.font.size = pptx.util.Pt(11)
            subtitle_text.font.bold = True
            subtitle_text.font.color.theme_color = 10
            subtitle_text.text = chunck_text(remove_newlines(v), MAX_CHARS_PER_LINE)
            # todo: change color
            bottom = subtitle.top + subtitle.height
        if k == 'text':
            # add text
            textbox_height = presentation.slide_height * textbox_height_estimate(v)
            if is_slide_full(presentation, bottom, textbox_height):
                # Create a new slide if the current slide is full
                slide_layout = presentation.slide_layouts[SLIDE_LAYOUT]
                slide = presentation.slides.add_slide(slide_layout)
                try:
                    banner_shape = next(find_template_shape(slide, BANNER_MARKUP))
                except:
                    banner_shape = slide.shapes[0]
                    banner_shape.text_frame.text = BANNER_MARKUP
                bottom = 0
            textbox = slide.shapes.add_textbox(0, bottom, presentation.slide_width, textbox_height)
            p = textbox.text_frame.paragraphs[0]
            p.font.size = pptx.util.Pt(10)
            p.text = chunck_text(remove_newlines(v), MAX_CHARS_PER_LINE)
            # p.alignment = PP_ALIGN.CENTER
            bottom = textbox.top + textbox.height
    return slide, bottom


def add_countries(presentation, slide, countries: List[Dict], top: int):
    """returns the bottom of the last item added"""
    MARGIN = 15_000

    bottom = top # at the beginning
    banner_shape = next(find_template_shape(slide, BANNER_MARKUP))
    for country in countries:
        # add banner with country name
        if is_slide_full(presentation, bottom + MARGIN, banner_shape.height):
            # Create a new slide if the current slide is full
            slide_layout = presentation.slide_layouts[SLIDE_LAYOUT]
            slide = presentation.slides.add_slide(slide_layout)
            try:
                banner_shape = next(find_template_shape(slide, BANNER_MARKUP))
            except:
                banner_shape = slide.shapes[0]
                banner_shape.text_frame.text = BANNER_MARKUP
            bottom = 0

        banner = clone_shape(banner_shape, banner_shape.left, bottom + MARGIN, banner_shape.width, BANNER_HEIGHT)
        banner.text_frame.text = country['title']
        banner.text_frame.paragraphs[0].font.color.theme_color = 14
        banner.text_frame.paragraphs[0].font.bold = True
        banner.text_frame.paragraphs[0].font.italic = True
        bottom = banner.top + banner.height + MARGIN

        for bodypart in country['body']:
            kv = tuple(bodypart.items())
            if len(kv) == 1:
                k, v = kv[0]
            else:
                raise Exception(f"Error while reading {bodypart} dict from article tree dict")
            
            if k == 'subtitle':
                # add subtitle
                subtitle_height = presentation.slide_height * textbox_height_estimate(v)
                subtitle = slide.shapes.add_textbox(SUBTITLE_LMARGIN, bottom, presentation.slide_width, subtitle_height)
                subtitle_text = subtitle.text_frame.paragraphs[0]
                subtitle_text.font.size = pptx.util.Pt(11)
                subtitle_text.font.bold = True
                subtitle_text.font.color.theme_color = 10
                subtitle_text.text = chunck_text(remove_newlines(v), MAX_CHARS_PER_LINE)
                # todo: change color
                bottom = subtitle.top + subtitle.height
            if k == 'text':
                # add text
                textbox_height = presentation.slide_height * textbox_height_estimate(v)
                if is_slide_full(presentation, bottom, textbox_height):
                    # Create a new slide if the current slide is full
                    slide_layout = presentation.slide_layouts[SLIDE_LAYOUT]
                    slide = presentation.slides.add_slide(slide_layout)
                    try:
                        banner_shape = next(find_template_shape(slide, BANNER_MARKUP))
                    except:
                        banner_shape = slide.shapes[0]
                        banner_shape.text_frame.text = BANNER_MARKUP
                    bottom = 0
                textbox = slide.shapes.add_textbox(0, bottom, presentation.slide_width, textbox_height)
                p = textbox.text_frame.paragraphs[0]
                p.font.size = pptx.util.Pt(10)
                p.text = chunck_text(remove_newlines(v), MAX_CHARS_PER_LINE)
                # p.alignment = PP_ALIGN.CENTER
                bottom = textbox.top + textbox.height
    return bottom


def add_sources(presentation, slide, sources: list):
    """add the sources and moves the slide to the end"""
    
    # add banner
    # slide_layout = presentation.slide_layouts[SLIDE_LAYOUT]
    # slide = presentation.slides.add_slide(slide_layout)
    banner_shape = list(find_template_shape(slide, SOURCE_BANNER_MARKUP))[0]
    banner = clone_shape(banner_shape, banner_shape.left, banner_shape.top, banner_shape.width, BANNER_HEIGHT)
    banner.height = int(BANNER_HEIGHT * 1.5)
    banner.text_frame.text = "Sources"
    banner.text_frame.paragraphs[0].font.color.theme_color = 14
    banner.text_frame.paragraphs[0].font.bold = True
    banner.text_frame.paragraphs[0].font.italic = True

    # add sources
    for i in range(len(sources)):
        template = list(find_template_shape(slide, SOURCE_MARKUP))[0]
        template.text_frame.paragraphs[i].text = sources[i]
        template.text_frame.paragraphs[i].runs[0].hyperlink.address = sources[i]
        template.text_frame.paragraphs[i].font.size = pptx.util.Pt(8)
    
    move_slide(presentation, presentation.slides.index(slide), len(list(presentation.slides)))

def add_credits(presentation, slide, author=None):
    """adds author name & moves the slide to the end"""
    if isinstance(author, str):
        for shape in find_template_shape(slide, AUTHOR_MARKUP):
            shape.text_frame.text = author
    move_slide(presentation, presentation.slides.index(slide), len(list(presentation.slides)))
    