import pptx
from word2pptx import (
    add_countries, 
    add_intro, 
    add_sources, 
    add_subtitle, 
    add_title, 
    clean_up_shapes, 
    word_tree,
    add_credits
)

state_dict = word_tree('powerpoints/veille_RPA.docx')


prs = pptx.Presentation("powerpoints/gabarit_v2.pptx")
top = 3_000_000 # decided arbitrarily
add_title(prs.slides[0], state_dict["title"])
add_subtitle(prs.slides[0], state_dict["subtitle"])
slide, bottom = add_intro(prs, prs.slides[0], state_dict["intro"], top)
bottom = add_countries(prs, slide, state_dict["countries"], bottom)
add_sources(prs, prs.slides[1], state_dict["sources"])
add_credits(prs, prs.slides[1], author=state_dict["author"]) # slide number is 1 since sources will be appended
# deletes the template shapes
clean_up_shapes(prs, "<banner>") 
clean_up_shapes(prs, "<source_banner>")
# clean_up_shapes(prs, "<source>") # issues with because it deletes the whole shape

prs.save('powerpoints/test.pptx')

