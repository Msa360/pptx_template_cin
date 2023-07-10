import pptx
from utils import add_countries, add_intro, add_sources, add_subtitle, add_title, clean_up_shapes
from docx_parser import full_tree

test_parsed_dict = {
    "title": "IoT",
    "subtitle": "Internet of things",
    "countries": [
        {
            "title": "France",
            "body": [
                {"subtitle": "Introduction"},
                {"text": "La France est un beau pays! "*14},
                {"subtitle": "Description"},
                {"text": "La France est vraiment un super beau pays! "*12}
            ]
        },
        {"title": "Australie", "body": [{"text": "L'allemagne est une³ beau pays! "*14}]},
        {"title": "Canada", "body": [{"text": "Le Canada de magne est un superbe de beau pays! "*14}]},
        {"title": "Kenya", "body": [{"text": "Le Kenya de magne est un vraiment beau pays! "*14}]},
        {"title": "Belgique", "body": [{"text": "La Belgique est un beau pays! j'y suis allé. "*14}]}
    ],
    "sources": [{"number": "1", "text": "www.example1.com"}, {"number": "2", "text": "www.example2.com/thisisanexample?iot=1"}]
}

state_dict = full_tree('powerpoints/edge_computing.docx')
# state_dict = full_tree('powerpoints/example.docx')

# prs = pptx.Presentation("powerpoints/test.pptx")
prs = pptx.Presentation("powerpoints/gabarit_mod.pptx")
top = 3000000
add_title(prs.slides[0], state_dict["title"])
add_subtitle(prs.slides[0], state_dict["subtitle"])
slide, bottom = add_intro(prs, prs.slides[0], state_dict["intro"], top)
bottom = add_countries(prs, slide, state_dict["countries"], bottom)
add_sources(prs, prs.slides[1], state_dict["sources"], 0)
# deletes the template shapes
clean_up_shapes(prs, "<banner>") 
clean_up_shapes(prs, "<source_banner>")
# clean_up_shapes(prs, "<source>") # issues with because it deletes the whole shape

prs.save('powerpoints/test2.pptx')

