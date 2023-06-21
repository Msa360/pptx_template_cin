import copy
import pptx
from pptx.shapes.autoshape import Shape
from pptx.enum.text import PP_ALIGN
from typing import Union, Optional

# GLOBAL VARS
SLIDE_LAYOUT = 4 # 6 is a blank one, 4 is good
MAX_CHARS_PER_LINE = 97 # approximate, experimentally determined

def chunck_text(text: str, lenght: int):
    """chuncks the text into lines that are smaller or equal than 'lenght' arg"""
    text = list(text)
    text_list = list(text)
    # for i in range(len(text) // lenght):  # iter over lines
    pos = 0
    while len(text_list) / lenght > 1:
        for j in reversed(range(0, lenght+1)):    # iter over chars
            if text_list[j] == " ":
                pos += j
                text[pos] = "\n"
                pos += 1
                text_list = text_list[j+1:]
                break
            if j == 0:
                raise Exception(f"Cannot chunck {''.join(text_list[:lenght])}|{text_list[lenght]}, it would break the word at '|'")
    return ''.join(text)

def remove_newlines(text: str):
    new_text = list(text)
    for i in range(len(text)):
        if 0 < i < len(text) - 1:
            if text[i] == '\n':
                if text[i-1] == ' ' and text[i+1] == ' ':
                    new_text[i] = ''
                    new_text[i-1] = ''
                elif text[i-1] == ' ' and text[i+1] != ' ':
                    new_text[i] = ''
                elif text[i-1] != ' ' and text[i+1] == ' ':
                    new_text[i] = ''
                elif text[i-1] != ' ' and text[i+1] != ' ':
                    new_text[i] = ' '
        else:
            if text[i] == '\n':
                new_text[i] = ''
    return new_text

# a very important algo
def textbox_height_estimate(text: str):
    """
    The goal is to estimate the height of the textbox to fit all the characters
    this is needed since pptx api can't do it because it would need a rendering engine
    """
    n_chars = len(text)
    return (max(n_chars - 100, 0)) * 0.00018 + 0.04


def is_slide_full(presentation, shape_top, shape_height):
    """check if the content will fit in the slide"""
    # Get the last slide in the presentation
    slide = presentation.slides[-1]
    shapes = list(slide.shapes)
    if len(shapes)==0 or presentation.slide_height > shape_top + shape_height:
        return False
    else:
        return True
    

def clone_shape(shape, left, top, width, height):
        # gotta be careful since the shape_id is read-only, it keeps the same id so the id isn't unique anymore
        # gotta be careful, the style can change because of hierarchies
        """
        Add a duplicate of `shape` to the same slide.
        Once duplicated, id is obselete, use `shape.element` to have uniqueness.
        """
        shape_obj = shape.element
        sp_tree = shape_obj.getparent()
        new_sp = copy.deepcopy(shape_obj)
        sp_tree.append(new_sp)
        new_shape = Shape(new_sp, None) #(GraphicFrame(new_sp, None) if isinstance(shape, GraphicFrame) else Shape(new_sp, None))
        new_shape.left = int(left)
        new_shape.top = int(top)
        new_shape.width = int(width)
        new_shape.height = int(height)
        return new_shape

def delete_shape(shape, slide):
    slide.shapes._spTree.remove(shape._element)

def find_template_shape(slide, markup: str):
    """markup is the template placeholder, raise error in markup isn't found"""
    used = 0
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.text == markup:
                yield shape
                used += 1
    if used == 0:
        raise Exception("template shape wasn't found")


def add_countries(presentation, slide, countries: dict, top: int):
    """returns the bottom of the last item added"""
    MARGIN = 15_000
    BANNER_HEIGHT = 400_000 # define arbitrarly to ressemble

    bottom = top # at the beginning
    banner_shape = next(find_template_shape(slide, "banner"))
    for country in countries:
        # add banner with country name
        if is_slide_full(presentation, bottom + MARGIN, banner_shape.height):
            print("full")
            # Create a new slide if the current slide is full
            slide_layout = presentation.slide_layouts[SLIDE_LAYOUT]
            slide = presentation.slides.add_slide(slide_layout)
            try:
                banner_shape = next(find_template_shape(slide, "banner"))
            except:
                banner_shape = slide.shapes[0]
            bottom = 0

        banner = clone_shape(banner_shape, banner_shape.left, bottom + MARGIN, banner_shape.width, BANNER_HEIGHT)
        banner.text_frame.text = country['title']
        banner.text_frame.paragraphs[0].font.color.theme_color = 14
        banner.text_frame.paragraphs[0].font.bold = True
        banner.text_frame.paragraphs[0].font.italic = True
        bottom = banner.top + banner.height + MARGIN

        for bodypart in country['body']:
            if bodypart['subtitle'] != None:
                # add subtitle
                subtitle_height = presentation.slide_height * textbox_height_estimate(bodypart['subtitle'])
                subtitle = slide.shapes.add_textbox(0, bottom, presentation.slide_width, subtitle_height)
                subtitle_text = subtitle.text_frame.paragraphs[0]
                subtitle_text.font.size = pptx.util.Pt(11)
                subtitle_text.font.bold = True
                subtitle_text.font.color.theme_color = 10
                subtitle_text.text = chunck_text(remove_newlines(bodypart['subtitle']), MAX_CHARS_PER_LINE)
                # todo: change color
                bottom = subtitle.top + subtitle.height

            # add text
            textbox_height = presentation.slide_height * textbox_height_estimate(bodypart['text'])
            if is_slide_full(presentation, bottom, textbox_height):
                print("full")
                # Create a new slide if the current slide is full
                slide_layout = presentation.slide_layouts[SLIDE_LAYOUT]
                slide = presentation.slides.add_slide(slide_layout)
                try:
                    banner_shape = next(find_template_shape(slide, "banner"))
                except:
                    banner_shape = slide.shapes[0]
                bottom = 0
            
            textbox = slide.shapes.add_textbox(0, bottom, presentation.slide_width, textbox_height)
            p = textbox.text_frame.paragraphs[0]
            p.font.size = pptx.util.Pt(10)
            p.text = chunck_text(remove_newlines(bodypart['text']), MAX_CHARS_PER_LINE)
            p.alignment = PP_ALIGN.CENTER
            bottom = textbox.top + textbox.height
    return bottom
